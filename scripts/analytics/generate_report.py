# -*- coding: utf-8 -*-
"""
エロナビ アクセス分析レポート 自動生成スクリプト

analytics/data/ 配下のJSONデータを読み込み、PowerPointレポートを自動生成する。
JSONがない場合はダミーデータでプレビュー生成。

Usage:
    python generate_report.py
    python generate_report.py --output /path/to/output.pptx
"""

import json
import os
import glob
import argparse
from datetime import datetime, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =============================================================================
# カラーパレット（ダークテーマ）
# =============================================================================
BG_COLOR = RGBColor(0x16, 0x1B, 0x22)
CARD_COLOR = RGBColor(0x1C, 0x23, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
ACCENT = RGBColor(0xE6, 0x39, 0x46)
GREEN = RGBColor(0x00, 0xCC, 0x66)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
BLUE = RGBColor(0x44, 0x88, 0xFF)
RED = RGBColor(0xFF, 0x44, 0x44)
DARK_ROW = RGBColor(0x1C, 0x23, 0x33)
DARKER_ROW = RGBColor(0x14, 0x19, 0x26)
HEADER_ROW = RGBColor(0x2A, 0x10, 0x14)

SLIDE_W_INCHES = 13.333
SLIDE_H_INCHES = 7.5

# =============================================================================
# データディレクトリ
# =============================================================================
SCRIPT_DIR = Path(__file__).resolve().parent
DATA_DIR = SCRIPT_DIR / "data"
OUTPUT_DIR = Path(r"c:\Users\atsus\000_ClaudeCode\051_FANZAアフィリ")


# =============================================================================
# ダミーデータ生成
# =============================================================================
def generate_dummy_data():
    """JSONが存在しない場合のダミーデータ"""
    today = datetime.now()

    # 日別推移データ（過去28日）
    daily_data = []
    for i in range(28, 0, -1):
        d = today - timedelta(days=i)
        import random
        random.seed(i * 7 + 42)
        pv = random.randint(80, 350)
        clicks = random.randint(10, 80)
        daily_data.append({
            "date": d.strftime("%Y-%m-%d"),
            "pv": pv,
            "clicks": clicks,
            "users": int(pv * 0.65),
        })

    # サマリー
    total_pv = sum(d["pv"] for d in daily_data)
    total_users = sum(d["users"] for d in daily_data)
    summary = {
        "period": f'{(today - timedelta(days=28)).strftime("%Y/%m/%d")} - {today.strftime("%Y/%m/%d")}',
        "total_pv": total_pv,
        "total_users": total_users,
        "total_articles": 287,
        "indexed_count": 245,
        "index_rate": round(245 / 287 * 100, 1),
        "avg_daily_pv": round(total_pv / 28),
        "data_available": False,
    }

    # 検索クエリTOP20
    queries = [
        ("エロナビ", 320, 4800, 6.7, 3.2),
        ("人妻 動画 おすすめ", 185, 9200, 2.0, 12.5),
        ("熟女 AV ランキング", 142, 7600, 1.9, 15.1),
        ("巨乳 新作", 128, 6100, 2.1, 8.7),
        ("素人 無修正", 115, 11000, 1.0, 22.3),
        ("VR AV おすすめ", 98, 4300, 2.3, 6.4),
        ("NTR 動画", 87, 3800, 2.3, 9.8),
        ("痴女 AV", 76, 3200, 2.4, 11.2),
        ("OL AV 人気", 68, 2900, 2.3, 14.6),
        ("レズ AV 新作", 62, 2700, 2.3, 10.3),
        ("コスプレ AV", 55, 2400, 2.3, 13.7),
        ("女子大生 AV", 51, 2200, 2.3, 16.2),
        ("ナンパ もの", 48, 2100, 2.3, 18.4),
        ("企画 AV 面白い", 42, 1800, 2.3, 20.1),
        ("アナル AV", 38, 1600, 2.4, 17.8),
        ("フェチ 動画", 35, 1500, 2.3, 19.5),
        ("3P AV", 32, 1400, 2.3, 21.3),
        ("潮吹き AV", 29, 1200, 2.4, 22.7),
        ("ロリ系 AV", 27, 1100, 2.5, 24.1),
        ("FANZA セール", 24, 900, 2.7, 5.3),
    ]
    search_queries = [
        {"query": q[0], "clicks": q[1], "impressions": q[2],
         "ctr": q[3], "position": q[4]} for q in queries
    ]

    # 人気ページTOP20
    pages = [
        ("/", 1820),
        ("/ranking/weekly", 890),
        ("/category/hitozuma", 675),
        ("/category/jukujo", 543),
        ("/category/kyonyu", 498),
        ("/ranking/monthly", 465),
        ("/category/vr", 412),
        ("/actress/popular", 387),
        ("/category/ntr", 356),
        ("/new-releases", 334),
        ("/category/chijo", 298),
        ("/category/ol", 267),
        ("/category/shiroto", 245),
        ("/category/cosplay", 223),
        ("/category/rezu", 198),
        ("/category/joshidaisei", 187),
        ("/sale", 176),
        ("/category/nanpa", 165),
        ("/category/anal", 154),
        ("/category/fetish", 143),
    ]
    popular_pages = [{"path": p[0], "pv": p[1]} for p in pages]

    # カテゴリ別アクセス
    categories = [
        ("人妻", 42, 675), ("熟女", 38, 543), ("巨乳", 35, 498),
        ("VR", 28, 412), ("NTR", 25, 356), ("痴女", 22, 298),
        ("OL", 18, 267), ("素人", 15, 245), ("コスプレ", 12, 223),
        ("レズ", 10, 198), ("女子大生", 8, 187), ("ナンパ", 7, 165),
        ("アナル", 6, 154), ("フェチ", 5, 143),
    ]
    category_data = [
        {"name": c[0], "articles": c[1], "pv": c[2]} for c in categories
    ]

    # 流入元分析
    traffic_sources = [
        {"source": "オーガニック検索", "sessions": 3200, "ratio": 52.1},
        {"source": "直接アクセス", "sessions": 1500, "ratio": 24.4},
        {"source": "SNS (Twitter/X)", "sessions": 890, "ratio": 14.5},
        {"source": "リファラル", "sessions": 420, "ratio": 6.8},
        {"source": "その他", "sessions": 130, "ratio": 2.1},
    ]

    # デバイス分析
    devices = [
        {"device": "モバイル", "sessions": 4200, "ratio": 68.3, "avg_duration": "1:42"},
        {"device": "PC", "sessions": 1650, "ratio": 26.8, "avg_duration": "3:15"},
        {"device": "タブレット", "sessions": 290, "ratio": 4.7, "avg_duration": "2:28"},
    ]

    # サイトヘルス
    site_health = {
        "valid_urls": 245,
        "error_urls": 12,
        "warning_urls": 30,
        "not_indexed": 42,
        "errors": [
            {"url": "/old-page-1", "type": "404"},
            {"url": "/old-page-2", "type": "404"},
            {"url": "/broken-link-3", "type": "404"},
            {"url": "/redirect-loop", "type": "リダイレクトエラー"},
            {"url": "/server-error-page", "type": "サーバーエラー"},
        ],
    }

    return {
        "summary": summary,
        "daily": daily_data,
        "search_queries": search_queries,
        "popular_pages": popular_pages,
        "categories": category_data,
        "traffic_sources": traffic_sources,
        "devices": devices,
        "site_health": site_health,
    }


# =============================================================================
# JSONデータ読み込み
# =============================================================================
def load_latest_json(pattern):
    """指定パターンに一致する最新のJSONファイルを読み込む"""
    files = sorted(glob.glob(str(DATA_DIR / pattern)), key=os.path.getmtime, reverse=True)
    if files:
        with open(files[0], "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def load_all_data():
    """全データを読み込み、ない場合はダミーで補完"""
    dummy = generate_dummy_data()

    summary = load_latest_json("summary*.json") or load_latest_json("ga4_summary*.json")
    daily = load_latest_json("daily*.json") or load_latest_json("ga4_daily*.json")
    queries = load_latest_json("search_queries*.json") or load_latest_json("sc_queries*.json")
    pages = load_latest_json("popular_pages*.json") or load_latest_json("ga4_pages*.json")
    categories = load_latest_json("categor*.json")
    traffic = load_latest_json("traffic*.json") or load_latest_json("ga4_traffic*.json")
    devices = load_latest_json("device*.json") or load_latest_json("ga4_device*.json")
    health = load_latest_json("site_health*.json") or load_latest_json("health*.json")

    data = {
        "summary": summary or dummy["summary"],
        "daily": daily or dummy["daily"],
        "search_queries": queries or dummy["search_queries"],
        "popular_pages": pages or dummy["popular_pages"],
        "categories": categories or dummy["categories"],
        "traffic_sources": traffic or dummy["traffic_sources"],
        "devices": devices or dummy["devices"],
        "site_health": health or dummy["site_health"],
    }

    # データ未取得フラグ
    data["has_real_data"] = any([
        summary, daily, queries, pages, categories, traffic, devices, health
    ])

    return data


# =============================================================================
# PowerPoint ヘルパー関数
# =============================================================================
class ReportBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W_INCHES)
        self.prs.slide_height = Inches(SLIDE_H_INCHES)
        self.slide_count = 0

    def add_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._set_bg(slide)
        self.slide_count += 1
        return slide

    def _set_bg(self, slide, color=BG_COLOR):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(self, slide, left, top, width, height, text,
                    font_size=14, color=LIGHT_GRAY, bold=False,
                    alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_multiline(self, slide, left, top, width, height, lines,
                      font_size=12, color=LIGHT_GRAY, line_spacing=1.3,
                      font_name="Meiryo"):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line_data in enumerate(lines):
            if isinstance(line_data, tuple):
                text, lcolor, lbold = line_data[0], line_data[1], line_data[2] if len(line_data) > 2 else False
            else:
                text, lcolor, lbold = str(line_data), color, False
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = lcolor
            p.font.bold = lbold
            p.font.name = font_name
            p.space_after = Pt(font_size * 0.4)
        return txBox

    def add_card(self, slide, left, top, width, height, fill_color=CARD_COLOR):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_title_bar(self, slide, text, subtitle=None):
        """スライド上部にアクセントバーとタイトル"""
        self.add_rect(slide, 0, 0, SLIDE_W_INCHES, 1.0, ACCENT)
        self.add_textbox(slide, 0.6, 0.12, 11, 0.7, text,
                         font_size=30, color=WHITE, bold=True)
        if subtitle:
            self.add_textbox(slide, 0.6, 0.62, 11, 0.35, subtitle,
                             font_size=13, color=RGBColor(0xFF, 0xCC, 0xCC))

    def add_page_number(self, slide, num, total):
        self.add_textbox(slide, 11.5, 7.0, 1.5, 0.4,
                         f"{num} / {total}", font_size=10,
                         color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    def add_table(self, slide, left, top, width, headers, rows,
                  col_widths=None, font_size=10, row_height=0.38):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(row_height * n_rows))
        tbl = tbl_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)

        # ヘッダー行
        for j, h in enumerate(headers):
            self._style_cell(tbl.cell(0, j), h, font_size=font_size,
                             color=WHITE, bg_color=HEADER_ROW, bold=True)

        # データ行
        for i, row in enumerate(rows):
            bg = DARK_ROW if i % 2 == 0 else DARKER_ROW
            for j, val in enumerate(row):
                self._style_cell(tbl.cell(i + 1, j), str(val),
                                 font_size=font_size, color=LIGHT_GRAY,
                                 bg_color=bg)

        return tbl_shape

    def _style_cell(self, cell, text, font_size=10, color=LIGHT_GRAY,
                    bg_color=None, bold=False, alignment=PP_ALIGN.CENTER):
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = "Meiryo"
            p.alignment = alignment
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_kpi_card(self, slide, left, top, width, height, label, value,
                     sub_text=None, value_color=WHITE):
        """KPI表示用のカード"""
        self.add_card(slide, left, top, width, height)
        self.add_textbox(slide, left + 0.2, top + 0.15, width - 0.4, 0.3,
                         label, font_size=12, color=MID_GRAY)
        self.add_textbox(slide, left + 0.2, top + 0.5, width - 0.4, 0.6,
                         str(value), font_size=32, color=value_color, bold=True)
        if sub_text:
            self.add_textbox(slide, left + 0.2, top + 1.1, width - 0.4, 0.3,
                             sub_text, font_size=10, color=MID_GRAY)

    def add_bar_chart_text(self, slide, left, top, width, height,
                           labels, values, max_val=None, bar_color=ACCENT):
        """テキストベースの横棒グラフ"""
        self.add_card(slide, left, top, width, height)
        if not values:
            return
        if max_val is None:
            max_val = max(values) if values else 1
        bar_area_width = width - 2.8
        row_h = min(0.32, (height - 0.4) / len(labels))

        for i, (label, val) in enumerate(zip(labels, values)):
            y = top + 0.2 + i * row_h
            # ラベル
            self.add_textbox(slide, left + 0.15, y, 1.2, row_h,
                             label, font_size=8, color=LIGHT_GRAY,
                             alignment=PP_ALIGN.RIGHT)
            # バー
            bar_w = max(0.05, (val / max_val) * bar_area_width) if max_val > 0 else 0.05
            self.add_rect(slide, left + 1.45, y + 0.04, bar_w, row_h - 0.08,
                          bar_color)
            # 値
            self.add_textbox(slide, left + 1.5 + bar_w, y, 1.0, row_h,
                             f"{val:,}", font_size=8, color=LIGHT_GRAY)

    def save(self, path):
        self.prs.save(path)
        print(f"レポート保存完了: {path}")


# =============================================================================
# スライド生成関数
# =============================================================================
def build_slide_title(rb, data):
    """スライド1: タイトル"""
    slide = rb.add_slide()
    today_str = datetime.now().strftime("%Y年%m月%d日")

    # 上部アクセントライン
    rb.add_rect(slide, 0, 0, SLIDE_W_INCHES, 0.06, ACCENT)
    # 下部アクセントライン
    rb.add_rect(slide, 0, SLIDE_H_INCHES - 0.06, SLIDE_W_INCHES, 0.06, ACCENT)

    # 左の装飾バー
    rb.add_rect(slide, 1.0, 1.5, 0.07, 4.0, ACCENT)

    rb.add_textbox(slide, 1.5, 2.0, 10, 1.2,
                   "エロナビ アクセス分析レポート",
                   font_size=46, color=WHITE, bold=True)
    rb.add_textbox(slide, 1.5, 3.4, 10, 0.6,
                   "PV / 検索クエリ / 流入元 / サイトヘルス 総合分析",
                   font_size=18, color=LIGHT_GRAY)
    rb.add_textbox(slide, 1.5, 4.3, 10, 0.5,
                   today_str, font_size=20, color=ACCENT, bold=True)

    if not data.get("has_real_data"):
        rb.add_card(slide, 1.5, 5.5, 5.0, 0.6, RGBColor(0x33, 0x20, 0x00))
        rb.add_textbox(slide, 1.7, 5.55, 4.6, 0.5,
                       "* データ未取得 - ダミーデータでプレビュー表示中",
                       font_size=12, color=YELLOW)

    # 装飾ボックス
    labels = ["PV分析", "検索クエリ", "流入元", "サイトヘルス"]
    for i, lbl in enumerate(labels):
        rb.add_card(slide, 8.0 + (i % 2) * 2.5, 2.0 + (i // 2) * 2.2, 2.2, 1.8)
        rb.add_textbox(slide, 8.0 + (i % 2) * 2.5 + 0.2, 2.6 + (i // 2) * 2.2,
                       1.8, 0.5, lbl, font_size=14, color=ACCENT, bold=True,
                       alignment=PP_ALIGN.CENTER)


def build_slide_summary(rb, data):
    """スライド2: サマリー"""
    slide = rb.add_slide()
    summary = data["summary"]
    rb.add_title_bar(slide, "サマリー", summary.get("period", ""))

    # KPIカード 4つ
    rb.add_kpi_card(slide, 0.5, 1.3, 2.9, 1.5,
                    "合計PV", f'{summary.get("total_pv", 0):,}',
                    f'日平均: {summary.get("avg_daily_pv", 0):,}',
                    value_color=ACCENT)
    rb.add_kpi_card(slide, 3.7, 1.3, 2.9, 1.5,
                    "ユーザー数", f'{summary.get("total_users", 0):,}',
                    value_color=BLUE)
    rb.add_kpi_card(slide, 6.9, 1.3, 2.9, 1.5,
                    "記事数", f'{summary.get("total_articles", 0):,}',
                    value_color=GREEN)
    rb.add_kpi_card(slide, 10.1, 1.3, 2.9, 1.5,
                    "インデックス率",
                    f'{summary.get("index_rate", 0)}%',
                    f'{summary.get("indexed_count", 0):,} / {summary.get("total_articles", 0):,}',
                    value_color=YELLOW)

    # 下部に追加情報
    rb.add_card(slide, 0.5, 3.2, 12.5, 3.8)
    info_lines = [
        ("主要KPI一覧", WHITE, True),
        ("", LIGHT_GRAY, False),
        (f'  合計PV: {summary.get("total_pv", 0):,}  |  '
         f'合計ユーザー: {summary.get("total_users", 0):,}  |  '
         f'日平均PV: {summary.get("avg_daily_pv", 0):,}', LIGHT_GRAY, False),
        (f'  記事数: {summary.get("total_articles", 0):,}  |  '
         f'インデックス済: {summary.get("indexed_count", 0):,}  |  '
         f'インデックス率: {summary.get("index_rate", 0)}%', LIGHT_GRAY, False),
        ("", LIGHT_GRAY, False),
        ("  ポイント:", ACCENT, True),
    ]

    # 自動ポイント生成
    index_rate = summary.get("index_rate", 0)
    if index_rate < 80:
        info_lines.append((f'    - インデックス率 {index_rate}% は改善余地あり。'
                           'noindexタグやクロールエラーを確認', YELLOW, False))
    elif index_rate >= 90:
        info_lines.append((f'    - インデックス率 {index_rate}% は良好。'
                           '引き続き新規記事のインデックス促進を', GREEN, False))
    else:
        info_lines.append((f'    - インデックス率 {index_rate}%。'
                           '未インデックスページの原因調査を推奨', LIGHT_GRAY, False))

    avg_pv = summary.get("avg_daily_pv", 0)
    if avg_pv < 100:
        info_lines.append(("    - 日平均PVが100未満。SEO・SNS施策の強化が必要", YELLOW, False))
    elif avg_pv >= 500:
        info_lines.append(("    - 日平均PV 500以上。順調に成長中", GREEN, False))

    rb.add_multiline(slide, 0.8, 3.4, 12.0, 3.4, info_lines, font_size=13)


def build_slide_daily_trend(rb, data):
    """スライド3: 日別推移グラフ"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "日別推移（過去28日）", "PV数・クリック数の日別推移")
    daily = data["daily"]

    if not daily:
        rb.add_textbox(slide, 3, 4, 7, 1, "データなし",
                       font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        return

    # PV棒グラフ（左半分）
    recent_14 = daily[-14:]
    pv_labels = [d.get("date", "")[-5:] for d in recent_14]
    pv_values = [d.get("pv", 0) for d in recent_14]
    rb.add_textbox(slide, 0.5, 1.2, 6, 0.4, "PV数（直近14日）",
                   font_size=14, color=WHITE, bold=True)
    rb.add_bar_chart_text(slide, 0.5, 1.6, 6.0, 5.2,
                          pv_labels, pv_values, bar_color=ACCENT)

    # クリック数棒グラフ（右半分）
    click_labels = [d.get("date", "")[-5:] for d in recent_14]
    click_values = [d.get("clicks", 0) for d in recent_14]
    rb.add_textbox(slide, 6.9, 1.2, 6, 0.4, "クリック数（直近14日）",
                   font_size=14, color=WHITE, bold=True)
    rb.add_bar_chart_text(slide, 6.9, 1.6, 6.0, 5.2,
                          click_labels, click_values, bar_color=BLUE)


def build_slide_search_queries(rb, data):
    """スライド4: 検索クエリTOP20"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "検索クエリ TOP20", "Search Console データ")
    queries = data["search_queries"][:20]

    headers = ["#", "検索クエリ", "クリック数", "表示回数", "CTR (%)", "平均順位"]
    rows = []
    for i, q in enumerate(queries, 1):
        rows.append([
            str(i),
            q.get("query", ""),
            f'{q.get("clicks", 0):,}',
            f'{q.get("impressions", 0):,}',
            f'{q.get("ctr", 0):.1f}',
            f'{q.get("position", 0):.1f}',
        ])

    col_widths = [0.4, 3.8, 1.3, 1.3, 1.0, 1.0]
    rb.add_table(slide, 0.5, 1.2, 12.5, headers, rows,
                 col_widths=col_widths, font_size=9, row_height=0.28)


def build_slide_popular_pages(rb, data):
    """スライド5: 人気ページTOP20"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "人気ページ TOP20", "PV数ベース")
    pages = data["popular_pages"][:20]

    headers = ["#", "ページパス", "PV数"]
    rows = []
    for i, p in enumerate(pages, 1):
        rows.append([
            str(i),
            p.get("path", ""),
            f'{p.get("pv", 0):,}',
        ])

    col_widths = [0.5, 8.5, 2.0]
    rb.add_table(slide, 0.5, 1.2, 12.5, headers, rows,
                 col_widths=col_widths, font_size=10, row_height=0.28)


def build_slide_categories(rb, data):
    """スライド6: カテゴリ別アクセス"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "カテゴリ別アクセス", "カテゴリごとの記事数とPV")
    categories = data["categories"]

    # テーブル（左側）
    headers = ["#", "カテゴリ", "記事数", "PV数"]
    rows = []
    for i, c in enumerate(categories, 1):
        rows.append([
            str(i),
            c.get("name", ""),
            str(c.get("articles", 0)),
            f'{c.get("pv", 0):,}',
        ])

    col_widths = [0.4, 2.5, 1.0, 1.2]
    rb.add_table(slide, 0.5, 1.2, 5.5, headers, rows,
                 col_widths=col_widths, font_size=10, row_height=0.34)

    # 右側に棒グラフ
    cat_names = [c.get("name", "") for c in categories[:10]]
    cat_pvs = [c.get("pv", 0) for c in categories[:10]]
    rb.add_textbox(slide, 6.5, 1.2, 6, 0.4, "PV数 TOP10",
                   font_size=14, color=WHITE, bold=True)
    rb.add_bar_chart_text(slide, 6.5, 1.6, 6.3, 5.2,
                          cat_names, cat_pvs, bar_color=GREEN)


def build_slide_traffic(rb, data):
    """スライド7: 流入元分析"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "流入元分析", "トラフィックソース内訳")
    sources = data["traffic_sources"]

    # テーブル
    headers = ["流入元", "セッション数", "割合 (%)"]
    rows = []
    for s in sources:
        rows.append([
            s.get("source", ""),
            f'{s.get("sessions", 0):,}',
            f'{s.get("ratio", 0):.1f}',
        ])

    col_widths = [3.0, 2.0, 1.5]
    rb.add_table(slide, 0.5, 1.3, 6.5, headers, rows,
                 col_widths=col_widths, font_size=12, row_height=0.45)

    # 右側にビジュアル表現（横棒）
    src_names = [s.get("source", "") for s in sources]
    src_ratios = [s.get("ratio", 0) for s in sources]
    colors = [ACCENT, BLUE, GREEN, YELLOW, MID_GRAY]

    rb.add_textbox(slide, 7.5, 1.3, 5, 0.4, "割合ビジュアル",
                   font_size=14, color=WHITE, bold=True)
    rb.add_card(slide, 7.5, 1.8, 5.3, 3.5)

    total = sum(s.get("sessions", 0) for s in sources)
    for i, s in enumerate(sources):
        y = 2.1 + i * 0.6
        ratio = s.get("ratio", 0)
        bar_w = max(0.1, ratio / 100 * 4.0)
        c = colors[i % len(colors)]
        rb.add_textbox(slide, 7.7, y - 0.15, 2.5, 0.3,
                       s.get("source", ""), font_size=10, color=LIGHT_GRAY)
        rb.add_rect(slide, 7.7, y + 0.12, bar_w, 0.18, c)
        rb.add_textbox(slide, 7.7 + bar_w + 0.1, y + 0.05, 1.5, 0.3,
                       f'{ratio:.1f}%', font_size=10, color=WHITE, bold=True)

    # 下部に分析コメント
    rb.add_card(slide, 0.5, 4.5, 12.5, 2.5)
    comment_lines = [("流入元分析コメント", WHITE, True)]
    organic = next((s for s in sources if "オーガニック" in s.get("source", "")), None)
    if organic:
        ratio = organic.get("ratio", 0)
        if ratio >= 50:
            comment_lines.append(
                (f"  - オーガニック検索が {ratio:.1f}% と主力。SEO施策が機能している",
                 GREEN, False))
        elif ratio < 30:
            comment_lines.append(
                (f"  - オーガニック検索 {ratio:.1f}% は低め。SEO強化でさらなる流入増が見込める",
                 YELLOW, False))

    sns = next((s for s in sources if "SNS" in s.get("source", "")), None)
    if sns:
        ratio = sns.get("ratio", 0)
        if ratio < 10:
            comment_lines.append(
                ("  - SNS流入が10%未満。Twitter/X投稿の頻度・質の改善を推奨", YELLOW, False))
        else:
            comment_lines.append(
                (f"  - SNS流入 {ratio:.1f}%。投稿戦略が一定の効果を発揮", GREEN, False))

    rb.add_multiline(slide, 0.8, 4.7, 12.0, 2.0, comment_lines, font_size=13)


def build_slide_devices(rb, data):
    """スライド8: デバイス分析"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "デバイス分析", "デバイス別セッション・滞在時間")
    devices = data["devices"]

    # KPIカード形式で表示
    device_colors = [BLUE, ACCENT, GREEN]
    for i, d in enumerate(devices):
        left = 0.5 + i * 4.2
        rb.add_card(slide, left, 1.3, 3.8, 2.5)
        rb.add_textbox(slide, left + 0.3, 1.5, 3.2, 0.4,
                       d.get("device", ""), font_size=18, color=WHITE, bold=True)
        c = device_colors[i % len(device_colors)]
        rb.add_textbox(slide, left + 0.3, 2.0, 3.2, 0.6,
                       f'{d.get("sessions", 0):,} セッション',
                       font_size=24, color=c, bold=True)
        rb.add_textbox(slide, left + 0.3, 2.6, 3.2, 0.4,
                       f'割合: {d.get("ratio", 0):.1f}%',
                       font_size=14, color=LIGHT_GRAY)
        avg_dur = d.get("avg_duration", "-")
        rb.add_textbox(slide, left + 0.3, 3.0, 3.2, 0.4,
                       f'平均滞在: {avg_dur}',
                       font_size=12, color=MID_GRAY)

    # テーブルも追加
    headers = ["デバイス", "セッション数", "割合 (%)", "平均滞在時間"]
    rows = []
    for d in devices:
        rows.append([
            d.get("device", ""),
            f'{d.get("sessions", 0):,}',
            f'{d.get("ratio", 0):.1f}',
            d.get("avg_duration", "-"),
        ])
    col_widths = [2.5, 2.0, 1.5, 2.0]
    rb.add_table(slide, 0.5, 4.3, 8.5, headers, rows,
                 col_widths=col_widths, font_size=11, row_height=0.45)

    # コメント
    mobile = next((d for d in devices if "モバイル" in d.get("device", "")), None)
    if mobile and mobile.get("ratio", 0) >= 60:
        rb.add_card(slide, 9.5, 4.3, 3.5, 1.8)
        rb.add_multiline(slide, 9.7, 4.5, 3.1, 1.4, [
            ("モバイル最適化", WHITE, True),
            (f'モバイル比率: {mobile.get("ratio", 0):.1f}%', ACCENT, False),
            ("モバイルUXの継続改善が重要", LIGHT_GRAY, False),
        ], font_size=11)


def build_slide_site_health(rb, data):
    """スライド9: サイトヘルス"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "サイトヘルス", "URL検査・エラー状況")
    health = data["site_health"]

    # KPI
    valid = health.get("valid_urls", 0)
    errors = health.get("error_urls", 0)
    warnings = health.get("warning_urls", 0)
    not_indexed = health.get("not_indexed", 0)

    rb.add_kpi_card(slide, 0.5, 1.3, 3.0, 1.5,
                    "正常URL", f'{valid:,}', value_color=GREEN)
    rb.add_kpi_card(slide, 3.8, 1.3, 3.0, 1.5,
                    "エラーURL", f'{errors:,}', value_color=RED)
    rb.add_kpi_card(slide, 7.1, 1.3, 3.0, 1.5,
                    "警告URL", f'{warnings:,}', value_color=YELLOW)
    rb.add_kpi_card(slide, 10.4, 1.3, 2.6, 1.5,
                    "未インデックス", f'{not_indexed:,}', value_color=MID_GRAY)

    # エラー一覧テーブル
    error_list = health.get("errors", [])
    if error_list:
        rb.add_textbox(slide, 0.5, 3.1, 6, 0.4, "エラー詳細一覧",
                       font_size=14, color=WHITE, bold=True)
        headers = ["#", "URL", "エラータイプ"]
        rows = []
        for i, e in enumerate(error_list[:15], 1):
            rows.append([
                str(i),
                e.get("url", ""),
                e.get("type", ""),
            ])
        col_widths = [0.4, 6.0, 2.0]
        rb.add_table(slide, 0.5, 3.5, 8.5, headers, rows,
                     col_widths=col_widths, font_size=10, row_height=0.32)

    # 右側にヘルス概要
    total_urls = valid + errors + warnings
    health_pct = round(valid / total_urls * 100, 1) if total_urls > 0 else 0
    rb.add_card(slide, 9.5, 3.1, 3.5, 3.5)
    health_color = GREEN if health_pct >= 90 else YELLOW if health_pct >= 70 else RED
    rb.add_multiline(slide, 9.7, 3.3, 3.1, 3.0, [
        ("ヘルススコア", WHITE, True),
        (f"{health_pct}%", health_color, True),
        ("", LIGHT_GRAY, False),
        (f"正常: {valid}", GREEN, False),
        (f"エラー: {errors}", RED, False),
        (f"警告: {warnings}", YELLOW, False),
        (f"未インデックス: {not_indexed}", MID_GRAY, False),
    ], font_size=12)


def build_slide_suggestions(rb, data):
    """スライド10: 改善提案"""
    slide = rb.add_slide()
    rb.add_title_bar(slide, "改善提案", "データに基づく自動分析")

    suggestions = generate_suggestions(data)

    # 提案を2列に分けて表示
    left_suggestions = suggestions[:len(suggestions) // 2 + 1]
    right_suggestions = suggestions[len(suggestions) // 2 + 1:]

    for col_idx, items in enumerate([left_suggestions, right_suggestions]):
        left = 0.5 + col_idx * 6.3
        for i, s in enumerate(items):
            y = 1.3 + i * 1.15
            priority_color = RED if s["priority"] == "高" else YELLOW if s["priority"] == "中" else GREEN
            rb.add_card(slide, left, y, 6.0, 1.0)
            # 優先度ラベル
            rb.add_rect(slide, left + 0.15, y + 0.15, 0.5, 0.3, priority_color)
            rb.add_textbox(slide, left + 0.15, y + 0.15, 0.5, 0.3,
                           s["priority"], font_size=10, color=WHITE, bold=True,
                           alignment=PP_ALIGN.CENTER)
            # タイトル
            rb.add_textbox(slide, left + 0.8, y + 0.1, 5.0, 0.35,
                           s["title"], font_size=13, color=WHITE, bold=True)
            # 詳細
            rb.add_textbox(slide, left + 0.8, y + 0.45, 5.0, 0.5,
                           s["detail"], font_size=10, color=LIGHT_GRAY)


def generate_suggestions(data):
    """データに基づいて改善提案を自動生成"""
    suggestions = []
    summary = data.get("summary", {})
    categories = data.get("categories", [])
    health = data.get("site_health", {})
    sources = data.get("traffic_sources", [])
    devices = data.get("devices", [])
    queries = data.get("search_queries", [])

    # インデックス率チェック
    index_rate = summary.get("index_rate", 100)
    if index_rate < 85:
        suggestions.append({
            "priority": "高",
            "title": "インデックス率の改善",
            "detail": f"現在 {index_rate}%。サイトマップ送信、内部リンク強化、"
                      "noindex誤設定の確認を推奨。",
        })

    # エラーURLチェック
    error_count = health.get("error_urls", 0)
    if error_count > 5:
        suggestions.append({
            "priority": "高",
            "title": "エラーURLの修正",
            "detail": f"{error_count}件のエラーURL。404リダイレクト設定、"
                      "壊れたリンクの修正が急務。",
        })

    # PV少ないカテゴリ
    if categories:
        low_pv_cats = [c for c in categories if c.get("pv", 0) < 200]
        if low_pv_cats:
            names = "、".join(c.get("name", "") for c in low_pv_cats[:3])
            suggestions.append({
                "priority": "中",
                "title": "低PVカテゴリへの注力",
                "detail": f"{names} 等のPVが低い。記事追加・内部リンク・"
                          "SNS露出で改善可能。",
            })

    # SNS流入チェック
    sns = next((s for s in sources if "SNS" in s.get("source", "")), None)
    if sns and sns.get("ratio", 0) < 15:
        suggestions.append({
            "priority": "中",
            "title": "SNS流入の強化",
            "detail": f"SNS流入 {sns.get('ratio', 0):.1f}%。定期的な"
                      "ポスト、トレンドハッシュタグ活用を推奨。",
        })

    # モバイル対応
    mobile = next((d for d in devices if "モバイル" in d.get("device", "")), None)
    if mobile and mobile.get("ratio", 0) >= 60:
        suggestions.append({
            "priority": "中",
            "title": "モバイルUXの最適化",
            "detail": f"モバイル比率 {mobile.get('ratio', 0):.1f}%。"
                      "ページ速度、タップしやすいUI、AMP対応を検討。",
        })

    # CTRの低いクエリ
    low_ctr = [q for q in queries if q.get("ctr", 0) < 2.0 and q.get("impressions", 0) > 3000]
    if low_ctr:
        names = "、".join(q.get("query", "") for q in low_ctr[:2])
        suggestions.append({
            "priority": "中",
            "title": "高表示・低CTRクエリの改善",
            "detail": f"「{names}」等、表示回数は多いがCTRが低い。"
                      "タイトル・ディスクリプションの最適化を。",
        })

    # 高順位キーワードの維持
    top_queries = [q for q in queries if q.get("position", 99) < 5]
    if top_queries:
        suggestions.append({
            "priority": "低",
            "title": "上位クエリの維持・強化",
            "detail": f"{len(top_queries)}件のクエリが5位以内。"
                      "コンテンツ更新・被リンク獲得で順位維持を。",
        })

    # 日平均PV
    avg_pv = summary.get("avg_daily_pv", 0)
    if avg_pv < 200:
        suggestions.append({
            "priority": "高",
            "title": "全体PVの底上げ",
            "detail": f"日平均PV {avg_pv:,}。ロングテールキーワード狙いの"
                      "記事量産、外部リンク獲得が効果的。",
        })

    # 最低でも3つ
    if len(suggestions) < 3:
        suggestions.append({
            "priority": "低",
            "title": "定期的なコンテンツ更新",
            "detail": "既存記事のリライト、新作情報の追加で鮮度を維持。"
                      "更新頻度を週2-3回以上に。",
        })

    return suggestions[:10]


# =============================================================================
# メイン処理
# =============================================================================
def generate_report(output_path=None):
    """レポートを生成して保存"""
    print("データを読み込み中...")
    data = load_all_data()

    if data.get("has_real_data"):
        print("実データを検出しました")
    else:
        print("JSONデータが見つかりません。ダミーデータでプレビュー生成します")

    print("レポートを生成中...")
    rb = ReportBuilder()

    # 全スライドを生成
    slide_builders = [
        build_slide_title,
        build_slide_summary,
        build_slide_daily_trend,
        build_slide_search_queries,
        build_slide_popular_pages,
        build_slide_categories,
        build_slide_traffic,
        build_slide_devices,
        build_slide_site_health,
        build_slide_suggestions,
    ]

    for builder in slide_builders:
        builder(rb, data)

    total = rb.slide_count
    for i, slide in enumerate(rb.prs.slides, 1):
        rb.add_page_number(slide, i, total)

    # 保存
    if output_path is None:
        date_str = datetime.now().strftime("%y%m%d")
        filename = f"エロナビアクセス分析レポート_{date_str}.pptx"
        output_path = str(OUTPUT_DIR / filename)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    rb.save(output_path)
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="エロナビ アクセス分析レポート自動生成")
    parser.add_argument("--output", "-o", type=str, default=None,
                        help="出力先パス（デフォルト: 日付付きファイル名）")
    args = parser.parse_args()
    generate_report(args.output)
